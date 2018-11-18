import csv
from pptx import Presentation

prs = Presentation()

slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Jakis teks"

tf = body_shape.text_frame
tf.text = "zawartość tekst frame"


with open ("report.csv")as f:
    data = csv.reader(f, deliiter=",")
    p = tf.add_paragraph()
    p.text = "kobiety"
    p.level = 1

p = tf.add_paragraph()
p.text = "przeżyło 70%"
p.level = 2


prs.save('przyklad.pptx')